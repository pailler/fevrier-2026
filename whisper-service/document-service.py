from flask import Flask, request, jsonify
import PyPDF2
from docx import Document
import requests
import os

app = Flask(__name__)

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'healthy'})

@app.route('/process', methods=['POST'])
def process_document():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        text = ""
        filename = file.filename.lower()
        
        if filename.endswith('.pdf'):
            # Process PDF
            pdf_reader = PyPDF2.PdfReader(file.stream)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        
        elif filename.endswith(('.docx', '.doc')):
            # Process Word document
            doc = Document(file.stream)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        
        elif filename.endswith(('.pptx', '.ppt')):
            # Process PowerPoint
            from pptx import Presentation
            prs = Presentation(file.stream)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        
        else:
            return jsonify({'error': 'Unsupported file format'}), 400
        
        return jsonify({'text': text.strip()})
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8080, debug=True)
